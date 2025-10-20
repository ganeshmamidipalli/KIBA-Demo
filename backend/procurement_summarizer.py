#!/usr/bin/env python3
"""
Procurement Scope Summarizer
----------------------------

Ingests uploaded project-scope files (PDF, DOCX, TXT, images, PPTX, XLSX/CSV, HTML),
extracts text, and calls an LLM (OpenAI Chat Completions API via HTTPS) to produce a
structured summary focused on procurement details.

Output: JSON per document with keys:
- document_path
- items: [
    {
      "product_name": str,
      "category": str,
      "budget": str | number | null,
      "quantity": str | number | null,
      "timeline": str | null,
      "notes": str,
      "summary": str
    }, ...
  ]
- overall_summary: str

Setup
-----
1) Python 3.10+
2) Install dependencies:
   pip install pypdf docx2txt pillow pytesseract python-pptx pandas lxml beautifulsoup4 requests openpyxl chardet
   # Optional fallback extractor
   pip install textract

   # For OCR: install Tesseract on your system.
   # macOS (brew): brew install tesseract
   # Ubuntu/Debian: sudo apt-get install tesseract-ocr
   # Windows: https://github.com/UB-Mannheim/tesseract/wiki

3) Set env vars (e.g. in .env or shell):
   export OPENAI_API_KEY="sk-..."
   export OPENAI_MODEL="gpt-4o-mini"   # or gpt-4o, gpt-4.1-mini, etc.

Usage
-----
python procurement_summarizer.py /path/to/uploads --out results.json
python procurement_summarizer.py file1.pdf file2.docx --print

"""
from __future__ import annotations
import os
import io
import re
import sys
import json
import glob
import math
import time
import base64
import argparse
import logging
from dataclasses import dataclass, asdict
from typing import List, Optional, Dict, Any, Iterable, Tuple

# Text extraction libs
from pypdf import PdfReader
try:
    import docx2txt
    DOCX2TXT_AVAILABLE = True
except ImportError:
    DOCX2TXT_AVAILABLE = False
    
try:
    from PIL import Image
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    
import pandas as pd
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False
    
import requests
try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

# Optional fallback
try:
    import textract  # type: ignore
    TEXTRACT_AVAILABLE = True
except Exception:
    TEXTRACT_AVAILABLE = False

logging.basicConfig(level=logging.INFO, format="[%(levelname)s] %(message)s")

# ---------------------------
# Utilities
# ---------------------------

def read_binary(path: str) -> bytes:
    with open(path, 'rb') as f:
        return f.read()


def sniff_decode(data: bytes) -> str:
    """Decode bytes to text with chardet fallback."""
    try:
        return data.decode('utf-8')
    except UnicodeDecodeError:
        if CHARDET_AVAILABLE:
            enc = chardet.detect(data).get('encoding') or 'utf-8'
            return data.decode(enc, errors='replace')
        return data.decode('utf-8', errors='replace')


# ---------------------------
# Extractors for each type
# ---------------------------

def extract_text_pdf(path: str) -> str:
    text_parts: List[str] = []
    with open(path, 'rb') as f:
        reader = PdfReader(f)
        for page in reader.pages:
            try:
                text_parts.append(page.extract_text() or "")
            except Exception as e:
                logging.warning(f"PDF page extract failed: {e}")
    return "\n".join(text_parts).strip()


def extract_text_docx(path: str) -> str:
    if not DOCX2TXT_AVAILABLE:
        return ""
    try:
        return docx2txt.process(path) or ""
    except Exception:
        return ""


def extract_text_txt(path: str) -> str:
    data = read_binary(path)
    return sniff_decode(data)


def extract_text_image(path: str) -> str:
    if not TESSERACT_AVAILABLE:
        logging.warning("PIL/pytesseract not available - cannot extract from images")
        return ""
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logging.error(f"OCR failed for {path}: {e}")
        return ""


def extract_text_pptx(path: str) -> str:
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""


def extract_text_spreadsheet(path: str) -> str:
    try:
        df_dict = pd.read_excel(path, sheet_name=None) if path.lower().endswith((".xlsx", ".xls")) else {"Sheet1": pd.read_csv(path)}
        parts = []
        for name, df in df_dict.items():
            parts.append(f"# Sheet: {name}\n")
            parts.append(df.to_csv(index=False))
        return "\n".join(parts)
    except Exception as e:
        logging.error(f"Spreadsheet parse failed for {path}: {e}")
        return ""


def extract_text_html(path: str) -> str:
    if not BS4_AVAILABLE:
        return ""
    try:
        html = read_binary(path).decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html, 'lxml')
        return soup.get_text(" ", strip=True)
    except Exception:
        return ""


def extract_text_generic(path: str) -> str:
    if TEXTRACT_AVAILABLE:
        try:
            text = textract.process(path)  # type: ignore
            return text.decode('utf-8', errors='ignore')
        except Exception as e:
            logging.warning(f"textract failed on {path}: {e}")
    return ""


EXTENSION_MAP = {
    ".pdf": extract_text_pdf,
    ".docx": extract_text_docx,
    ".txt": extract_text_txt,
    ".md": extract_text_txt,
    ".rtf": extract_text_generic,
    ".png": extract_text_image,
    ".jpg": extract_text_image,
    ".jpeg": extract_text_image,
    ".tif": extract_text_image,
    ".tiff": extract_text_image,
    ".bmp": extract_text_image,
    ".pptx": extract_text_pptx,
    ".xlsx": extract_text_spreadsheet,
    ".xls": extract_text_spreadsheet,
    ".csv": extract_text_spreadsheet,
    ".html": extract_text_html,
    ".htm": extract_text_html,
}


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    extractor = EXTENSION_MAP.get(ext)
    if extractor is not None:
        text = extractor(path)
    else:
        text = extract_text_generic(path)
    # Clean up huge whitespace and limit extremely long docs (to control token cost)
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) > 40000:
        logging.info(f"Truncating long text for {path} to 40k characters.")
        text = text[:40000]
    return text


# ---------------------------
# OpenAI Chat Completions via HTTPS
# ---------------------------

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
OPENAI_BASE = os.getenv("OPENAI_BASE", "https://api.openai.com/v1")

if not OPENAI_API_KEY:
    logging.warning("OPENAI_API_KEY is not set. Set it before running to enable LLM summarization.")


def llm_extract_procurement(text: str, model: str = OPENAI_MODEL) -> Dict[str, Any]:
    """Ask the LLM to extract structured procurement info from free text.
    Uses Chat Completions with JSON mode for consistent parsing.
    """
    url = f"{OPENAI_BASE}/chat/completions"

    system = (
        "You are a precise procurement analyst. Extract procurement-relevant details from scope documents."
        " Always return STRICT JSON that matches the provided schema."
    )

    schema_hint = {
        "type": "object",
        "properties": {
            "items": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "product_name": {"type": ["string", "null"]},
                        "category": {"type": ["string", "null"]},
                        "budget": {"type": ["string", "number", "null"]},
                        "quantity": {"type": ["string", "number", "null"]},
                        "timeline": {"type": ["string", "null"]},
                        "notes": {"type": ["string", "null"]},
                        "summary": {"type": ["string", "null"]},
                    },
                    "required": ["summary"],
                    "additionalProperties": False,
                },
            },
            "overall_summary": {"type": "string"},
        },
        "required": ["items", "overall_summary"],
        "additionalProperties": False,
    }

    user = (
        "Extract procurement details. Focus on product name, category, budget, quantity,"
        " timeline/milestones, and any constraints/specs. If multiple products are mentioned,"
        " return multiple items. Keep summaries concise but complete."
        f"\n\nDOCUMENT:\n{text}"
    )

    payload = {
        "model": model,
        "response_format": {"type": "json_object"},
        "messages": [
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        "temperature": 0.0,
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {OPENAI_API_KEY}",
    }

    try:
        r = requests.post(url, headers=headers, json=payload, timeout=120)
        r.raise_for_status()
        data = r.json()
        content = data["choices"][0]["message"]["content"]
        logging.debug(f"LLM response content: {content[:500]}")
        try:
            result = json.loads(content)
            logging.debug(f"Parsed JSON items: {len(result.get('items', []))}")
            return result
        except json.JSONDecodeError as e:
            logging.error(f"JSON decode error: {e}")
            # Best-effort fallback: wrap in a generic container
            return {"items": [{"summary": content}], "overall_summary": content}
    except Exception as e:
        logging.error(f"LLM extraction failed: {e}")
        return {"items": [{"summary": f"Error: {e}"}], "overall_summary": f"Error: {e}"}


# ---------------------------
# Core pipeline
# ---------------------------

@dataclass
class DocResult:
    document_path: str
    items: List[Dict[str, Any]]
    overall_summary: str


def process_path(path: str) -> DocResult:
    text = extract_text(path)
    if not text:
        return DocResult(document_path=path, items=[{"summary": "No text could be extracted."}], overall_summary="No extractable text.")

    if not OPENAI_API_KEY:
        # Offline mode: simple heuristic stub
        return DocResult(
            document_path=path,
            items=[{"product_name": None, "category": None, "budget": None, "quantity": None, "timeline": None, "notes": None, "summary": text[:500] + ("..." if len(text) > 500 else "")}],
            overall_summary="Set OPENAI_API_KEY to enable LLM-based structured extraction.",
        )

    result = llm_extract_procurement(text)
    # Normalize keys and types
    items = result.get("items") or []
    if not isinstance(items, list):
        items = [items]
    overall = result.get("overall_summary") or ""
    return DocResult(document_path=path, items=items, overall_summary=overall)


def iter_input_paths(args: argparse.Namespace) -> Iterable[str]:
    if args.paths:
        for p in args.paths:
            if os.path.isdir(p):
                for ext in EXTENSION_MAP.keys():
                    yield from glob.glob(os.path.join(p, f"**/*{ext}"), recursive=True)
            else:
                yield p


def main():
    parser = argparse.ArgumentParser(description="Summarize procurement scope from uploaded files.")
    parser.add_argument('paths', nargs='+', help='Files or directories to process')
    parser.add_argument('--out', help='Write all results to a JSON file')
    parser.add_argument('--print', dest='do_print', action='store_true', help='Print results to stdout')
    parser.add_argument('--model', default=os.getenv('OPENAI_MODEL', OPENAI_MODEL), help='Override model name')
    args = parser.parse_args()

    results: List[Dict[str, Any]] = []
    for path in iter_input_paths(args):
        logging.info(f"Processing {path}")
        try:
            res = process_path(path)
            results.append(asdict(res))
        except Exception as e:
            logging.exception(f"Failed on {path}: {e}")
            results.append({
                "document_path": path,
                "items": [{"summary": f"Error: {e}"}],
                "overall_summary": "",
            })

    if args.do_print or not args.out:
        print(json.dumps(results, indent=2, ensure_ascii=False))

    if args.out:
        with open(args.out, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2, ensure_ascii=False)
        logging.info(f"Wrote {args.out}")


if __name__ == '__main__':
    main()

